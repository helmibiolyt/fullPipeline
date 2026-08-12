#!/usr/bin/env python3
"""Build the investor deck.

    python presentation/make_deck.py

Roughly sixty slides covering architecture, every source in use, how the graph
is joined, why certain entities are embedded, the vector store, orchestration,
and the hardware limits. Figures come from content.py, which reads the repo and
the lake inventory, so regenerating after a change updates the deck.
"""
from __future__ import annotations

import pathlib
import sys

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = pathlib.Path(__file__).resolve().parent
sys.path.insert(0, str(HERE))

import content as C          # noqa: E402
import theme as T            # noqa: E402

OUT = HERE / "Biolyt_Platform.pptx"
FULL = Inches(11.9)
_n = [0]


# --------------------------------------------------------------------------
def new(prs, kicker=None, title=None, color=T.BLUE, page=True):
    s = T.blank(prs)
    T.bg(s)
    y = T.slide_title(s, title, kicker, color) if title else Inches(0.7)
    _n[0] += 1
    if page:
        T.footer(s, _n[0])
    T.transition(s, "fade")
    return s, y


def section(prs, num, title, sub, color=T.BLUE):
    s = T.blank(prs)
    T.bg(s)
    _n[0] += 1
    T.solid(T.rect(s, 0, 0, Inches(0.14), T.H), color)
    a = T.text(s, T.MARGIN, Inches(2.5), Inches(3), Inches(1.1),
               f"{num:02d}", size=76, color=color, bold=True)
    b = T.text(s, T.MARGIN, Inches(3.5), Inches(10), Inches(0.9), title,
               size=42, color=T.TEXT, bold=True)
    c = T.text(s, T.MARGIN, Inches(4.5), Inches(9), Inches(0.8), sub,
               size=15, color=T.MUTED, line=1.35)
    T.footer(s, _n[0])
    T.transition(s, "push")
    T.animate(s, [a, b, c], step=140)
    return s


def bullets(items, size=13.5, color=None, gap=7):
    out = []
    for it in items:
        if isinstance(it, tuple):
            head, body = it
            out.append((head, {"size": size + 0.5, "bold": True,
                               "color": color or T.TEXT, "space_after": 1}))
            out.append((body, {"size": size, "color": T.MUTED,
                               "space_after": gap}))
        else:
            out.append((it, {"size": size, "color": color or T.MUTED,
                             "space_after": gap}))
    return out


def source_card(s, x, y, w, h, d):
    """One data source: what it is, what it builds, real columns, real values."""
    T.rect(s, x, y, w, h, T.PANEL, radius=0.05)
    accent = {"graph": T.BLUE, "docs": T.VIOLET,
              "both": T.CYAN}[d["role"]]
    T.solid(T.rect(s, x, y, Pt(3.5), h), accent)

    pad = Inches(0.26)
    lw = w * 0.44
    T.text(s, x + pad, y + Inches(0.14), lw - pad, Inches(0.3),
           d["src"], size=14.5, color=T.TEXT, bold=True)

    cx = x + pad
    label = {"graph": "GRAPH", "docs": "DOCUMENTS", "both": "GRAPH + DOCS"}
    cx = T.chip(s, cx, y + Inches(0.46), label[d["role"]], accent)
    vol = (f"{d['pdf']:,} PDF" if d["role"] == "docs"
           else f"{d['n_csv']} CSV · {C.human(d['csv_bytes'])}")
    T.text(s, cx + Inches(0.04), y + Inches(0.49), Inches(2.2), Inches(0.24),
           vol, size=10, color=T.DIM)

    purpose = d["purpose"]
    if len(purpose) > 128:
        cut = purpose.rfind(" ", 0, 128)
        purpose = purpose[:cut if cut > 90 else 128].rstrip(" ,.") + "…"
    T.text(s, x + pad, y + Inches(0.80), lw - pad * 1.4, Inches(0.44),
           purpose, size=10.5, color=T.MUTED, line=1.24)
    if d["builds"]:
        b = ", ".join(v.split(":")[-1] for v in d["builds"][:5])
        T.text(s, x + pad, y + h - Inches(0.3), lw - pad * 1.2, Inches(0.24),
               "BUILDS  " + b[:58], size=9, color=T.DIM, bold=True)

    # Right: real columns and the first real row's values.
    rx = x + lw
    cols = C.featured_columns(d, n=3)
    if cols:
        T.text(s, rx, y + Inches(0.16), Inches(3.4), Inches(0.22),
               "KEY COLUMNS", size=9, color=T.DIM, bold=True)
        T.text(s, rx + Inches(2.55), y + Inches(0.16), Inches(4), Inches(0.22),
               "EXAMPLE VALUE", size=9, color=T.DIM, bold=True)
        for i, (col, val) in enumerate(cols):
            yy = y + Inches(0.42 + i * 0.3)
            T.text(s, rx, yy, Inches(2.5), Inches(0.26), col,
                   size=10, color=T.TEXT, font=T.FONT_MONO)
            T.text(s, rx + Inches(2.55), yy, w - lw - Inches(2.75),
                   Inches(0.26), val, size=10, color=T.MUTED,
                   font=T.FONT_MONO)
    elif d["role"] == "docs":
        T.text(s, rx, y + Inches(0.4), w - lw - pad, Inches(0.8),
               f"{d['pdf']:,} PDF documents · {C.human(d['doc_bytes'])}\n"
               "No CSV — this source feeds the vector store only.",
               size=11, color=T.MUTED, line=1.3)



# ==========================================================================
def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = T.W, T.H
    U = C.used_sources()

    # ---------------------------------------------------------------- 1
    s = T.blank(prs)
    T.bg(s)
    T.solid(T.rect(s, 0, 0, T.W, Inches(0.09)), T.BLUE)
    a = T.text(s, T.MARGIN, Inches(2.15), Inches(11), Inches(0.4),
               "BIOMEDICAL INTELLIGENCE INFRASTRUCTURE", size=14,
               color=T.BLUE, bold=True)
    b = T.text(s, T.MARGIN, Inches(2.66), Inches(11.5), Inches(1.5),
               "One connected view of\nglobal drug development",
               size=52, color=T.TEXT, bold=True, line=1.08)
    c = T.text(s, T.MARGIN, Inches(4.5), Inches(9.6), Inches(0.9),
               f"50 data sources, reconciled into a "
               f"{C.millions(C.live_counts()['_nodes'])}-node knowledge graph "
               "and a 3.2M-chunk document store — kept current automatically.",
               size=15.5, color=T.MUTED, line=1.4)
    stats = T.rect(s, T.MARGIN, Inches(5.65), FULL, Inches(1.05),
                   T.PANEL, radius=0.06)
    for i, (v, l, col) in enumerate([
            (C.millions(C.live_counts()["_nodes"]), "graph nodes", T.BLUE),
            (C.millions(C.live_counts()["_edges"]), "relationships", T.BLUE),
            ("3.24M", "document chunks", T.VIOLET),
            ("93,505", "documents", T.VIOLET),
            ("41", "live sources", T.GREEN)]):
        T.stat(s, T.MARGIN + Inches(0.35 + i * 2.32), Inches(5.82),
               Inches(2.2), v, l, col, size=27, lsize=9.5)
    _n[0] += 1
    T.transition(s, "fade")
    T.animate(s, [a, b, c, stats], step=200)

    # ---------------------------------------------------------------- 2
    s, y = new(prs, "the short version", "What we built")
    cards = [
        T.card(s, T.MARGIN, y, Inches(3.83), Inches(2.0), "Collect",
               "50 scrapers pull from registries, agencies, ontologies and "
               "literature into one S3 lake. Each publishes a full snapshot "
               "on a schedule.", T.CYAN),
        T.card(s, T.MARGIN + Inches(4.03), y, Inches(3.83), Inches(2.0),
               "Connect",
               "35 sources are reconciled into a knowledge graph — 22 entity "
               "types, 32 relationship types — resolving names that no shared "
               "identifier links.", T.BLUE),
        T.card(s, T.MARGIN + Inches(8.06), y, Inches(3.84), Inches(2.0),
               "Retrieve",
               "93,505 regulatory documents chunked and embedded, searchable "
               "by meaning and joined back to the graph by identifier.",
               T.VIOLET),
    ]
    band = T.rect(s, T.MARGIN, y + Inches(2.25), FULL, Inches(1.5),
                  T.PANEL_2, radius=0.05)
    T.text(s, T.MARGIN + Inches(0.34), y + Inches(2.45), FULL - Inches(0.7),
           Inches(1.1),
           [("The hard part is not collecting the data. It is joining it.",
             {"size": 17, "color": T.TEXT, "bold": True, "space_after": 6}),
            ("The FDA writes ATORVASTATIN CALCIUM, Health Canada writes "
             "Atorvastatin calcium (as trihydrate), a trial registry writes "
             "Lipitor 40mg, ChEMBL writes CHEMBL1487. No identifier is common "
             "to all four. Everything we built rests on resolving that.",
             {"size": 12.5, "color": T.MUTED})])
    T.animate(s, cards + [band], step=180)

    # ---------------------------------------------------------------- 3
    s, y = new(prs, "contents", "What this deck covers")
    items = [("01  Architecture", "Two hosts, one lake, and why they are apart"),
             ("02  Data acquisition", "41 live sources, what each contributes"),
             ("03  Building the graph", "Normalisation, resolution, entities, edges"),
             ("04  Embedding entities", "Why some concepts need vectors"),
             ("05  The document store", "Chunking, embedding, retrieval"),
             ("06  Orchestration", "How everything stays current, automatically"),
             ("07  Limits and roadmap", "What the hardware costs us today")]
    shapes = []
    for i, (h, b) in enumerate(items):
        col = i % 2
        row = i // 2
        x = T.MARGIN + col * Inches(6.05)
        yy = y + Inches(row * 0.86)
        card = T.rect(s, x, yy, Inches(5.85), Inches(0.72), T.PANEL,
                      radius=0.08)
        T.text(s, x + Inches(0.28), yy + Inches(0.1), Inches(5.4),
               Inches(0.28), h, size=14, color=T.TEXT, bold=True)
        T.text(s, x + Inches(0.28), yy + Inches(0.38), Inches(5.4),
               Inches(0.26), b, size=11, color=T.MUTED)
        shapes.append(card)
    T.animate(s, shapes, step=110)

    # ================================================== 01 ARCHITECTURE
    section(prs, 1, "Architecture",
            "One lake, two stores, and a scheduler that keeps them in step "
            "without anyone watching.")

    s, y = new(prs, "the problem", "Forty-nine sources that were never meant to join")
    left = T.rect(s, T.MARGIN, y, Inches(5.85), Inches(3.5), T.PANEL,
                  radius=0.05)
    T.text(s, T.MARGIN + Inches(0.3), y + Inches(0.22), Inches(5.3),
           Inches(0.3), "The same drug, four ways", size=15, color=T.TEXT,
           bold=True)
    T.code(s, T.MARGIN + Inches(0.3), y + Inches(0.68), Inches(5.25),
           Inches(1.55),
           "FDA          ATORVASTATIN CALCIUM\n"
           "Health Canada Atorvastatin calcium\n"
           "             (as trihydrate)\n"
           "Trial reg.    Lipitor 40mg\n"
           "ChEMBL        CHEMBL1487")
    T.text(s, T.MARGIN + Inches(0.3), y + Inches(2.4), Inches(5.3),
           Inches(0.9),
           "No identifier is common to all four. A join on name fails; a join "
           "on any single id covers a fraction of the data.",
           size=12, color=T.MUTED, line=1.3)

    right = T.rect(s, T.MARGIN + Inches(6.05), y, Inches(5.85), Inches(3.5),
                   T.PANEL_2, radius=0.05)
    T.text(s, T.MARGIN + Inches(6.35), y + Inches(0.22), Inches(5.3),
           Inches(0.3), "What that costs the industry", size=15,
           color=T.TEXT, bold=True)
    T.text(s, T.MARGIN + Inches(6.35), y + Inches(0.7), Inches(5.25),
           Inches(2.6),
           bullets([
               ("Analysts reconcile by hand",
                "Competitive and regulatory questions are answered by opening "
                "six portals and matching names by eye."),
               ("Questions that span sources go unasked",
                "\"Which drugs targeting this protein were tested in the Gulf "
                "and carry a cardiac safety signal\" spans four registries."),
               ("Answers cannot be audited",
                "A spreadsheet does not record how two rows were matched.")]),
           line=1.28)
    T.animate(s, [left, right], step=200)

    s, y = new(prs, "system architecture", "How the platform fits together")
    _architecture(s, y)

    s, y = new(prs, "infrastructure", "Two hosts, and why they cannot be one")
    T.table(s, T.MARGIN, y, FULL,
            ["", "Graph host (Azure)", "Vector host (AWS)"],
            [["Runs", "Neo4j + nightly graph build", "Qdrant + ingest + search API + Airflow"],
             ["Memory profile", "4 GB heap + 4 GB page cache, 6 GB build peak",
              "~8 GB of vectors held resident"],
             ["Irreplaceable state", "None — everything derives from S3",
              "~27 GB of embeddings"],
             ["If lost", "Rebuild in ~25 min, automatic",
              "Hours on GPU, days on CPU"],
             ["Strategy", "Throw away and rebuild", "Back up and migrate"]],
            col_w=[2.0, 4.0, 4.0], size=11.5, row_h=Inches(0.44))
    note = T.rect(s, T.MARGIN, y + Inches(2.9), FULL, Inches(1.25),
                  T.PANEL_2, radius=0.05)
    T.solid(T.rect(s, T.MARGIN, y + Inches(2.9), Pt(3.5), Inches(1.25)), T.AMBER)
    T.text(s, T.MARGIN + Inches(0.32), y + Inches(3.08), FULL - Inches(0.7),
           Inches(0.95),
           [("They are apart because their memory profiles conflict.",
             {"size": 14, "color": T.TEXT, "bold": True, "space_after": 4}),
            ("Qdrant wants its vectors resident and Neo4j wants its store "
             "resident. Two of those on one 16 GB box means whichever is "
             "larger gets paged out — and a swapped page cache is slower than "
             "a small one. This is the constraint that shapes the whole "
             "deployment.", {"size": 12, "color": T.MUTED})])
    T.animate(s, [note], step=150)

    # ================================================== 02 SOURCES
    section(prs, 2, "Data acquisition",
            "Forty-one sources in production across eight domains — each with "
            "a defined purpose, not a bulk crawl.", T.CYAN)

    s, y = new(prs, "anatomy", "What a scraper is", T.CYAN)
    T.code(s, T.MARGIN, y, Inches(5.6), Inches(1.75),
           "scrape/<Category>/<site>/\n"
           "    manifest.yaml    name, s3_base, enabled,\n"
           "                     schedule, args, timeout\n"
           "    scrape.py        fetch, normalise, write")
    T.text(s, T.MARGIN, y + Inches(2.0), Inches(5.6), Inches(1.6),
           bullets([
               ("The manifest is the contract",
                "The runner discovers sources from the folder tree and reads "
                "the manifest. No source is hardcoded anywhere."),
               ("enabled is the switch",
                "A disabled source is not fetched, publishes nothing, and so "
                "triggers nothing downstream.")]), line=1.3)
    cards = []
    for i, (t, b, col) in enumerate([
            ("Adding a source is additive",
             "A new folder plus a manifest. Nothing else in the system "
             "changes — the graph declares which files it reads separately.",
             T.CYAN),
            ("Failure is isolated",
             "One source failing does not stop the others, and does not "
             "corrupt what is already published.", T.BLUE),
            ("Cost is proportional",
             "Sources run on their own schedule. A source nobody needs this "
             "week costs nothing.", T.GREEN)]):
        cards.append(T.card(s, T.MARGIN + Inches(6.05), y + Inches(i * 1.28),
                            Inches(5.85), Inches(1.15), t, b, col,
                            title_size=14, body_size=11.5))
    T.animate(s, cards, step=170)

    s, y = new(prs, "pipeline stages", "What one run does", T.CYAN)
    _stages(s, y)

    s, y = new(prs, "why it never double counts",
               "A commit replaces — it does not append", T.CYAN)
    T.text(s, T.MARGIN, y, Inches(6.0), Inches(2.4),
           bullets([
               ("Each source publishes a full snapshot",
                "One stable key per file, overwritten every run. "
                "ClinicalTrials.gov holds exactly two objects: a 2.9 GB CSV "
                "and _LATEST.json."),
               ("The graph is rebuilt, not patched",
                "Every sync rebuilds from current snapshots and replaces the "
                "store. There is no incremental append to get wrong."),
               ("Proven, not asserted",
                "Two consecutive builds from the same snapshot produced "
                "identical counts — 1,049,701 trials both times.")]), line=1.3)
    proof = T.rect(s, T.MARGIN + Inches(6.4), y, Inches(5.5), Inches(2.75),
                   T.PANEL, radius=0.05)
    T.text(s, T.MARGIN + Inches(6.7), y + Inches(0.2), Inches(5), Inches(0.3),
           "Live sync test · 30 July 2026", size=13, color=T.GREEN, bold=True)
    T.table(s, T.MARGIN + Inches(6.7), y + Inches(0.62), Inches(4.9),
            ["ClinicalTrial nodes", "before", "after"],
            [["registry = ctgov", "595,630", "596,490"],
             ["all registries", "1,048,841", "1,049,701"],
             ["net change", "—", "+860"]],
            col_w=[2.4, 1.5, 1.5], size=11, head_color=T.GREEN)
    T.text(s, T.MARGIN + Inches(6.7), y + Inches(2.25), Inches(5),
           Inches(0.4),
           "Enabling one scraper drove scrape → build → validate → import "
           "with no manual step.", size=10.5, color=T.MUTED, line=1.25)
    T.animate(s, [proof], step=150)

    s, y = new(prs, "coverage", "Eight domains", T.CYAN)
    shapes = []
    for i, cat in enumerate(C.CATEGORY_ORDER):
        lst = U.get(cat, [])
        col, row = i % 4, i // 4
        x = T.MARGIN + col * Inches(3.0)
        yy = y + row * Inches(1.72)
        cd = T.rect(s, x, yy, Inches(2.82), Inches(1.55), T.PANEL, radius=0.06)
        T.solid(T.rect(s, x, yy, Inches(2.82), Pt(3)), T.CYAN)
        T.text(s, x + Inches(0.22), yy + Inches(0.18), Inches(2.4),
               Inches(0.5), C.CATEGORY_LABEL[cat], size=12.5, color=T.TEXT,
               bold=True, line=1.15)
        T.text(s, x + Inches(0.22), yy + Inches(0.76), Inches(2.4),
               Inches(0.7), C.CATEGORY_WHY[cat][:112], size=9.5,
               color=T.MUTED, line=1.25)
        T.text(s, x + Inches(0.22), yy + Inches(1.28), Inches(2.4),
               Inches(0.24), f"{len(lst)} live sources", size=10,
               color=T.CYAN, bold=True)
        shapes.append(cd)
    T.animate(s, shapes, step=90)

    # Per-category source slides
    for cat in C.CATEGORY_ORDER:
        lst = U.get(cat, [])
        for chunk_i in range(0, len(lst), 3):
            part = lst[chunk_i:chunk_i + 3]
            more = f"  ({chunk_i // 3 + 1}/{(len(lst) + 2) // 3})" \
                if len(lst) > 3 else ""
            s, y = new(prs, C.CATEGORY_LABEL[cat] + more,
                       _cat_title(cat, chunk_i), T.CYAN)
            shapes = []
            for j, d in enumerate(part):
                yy = y + Inches(j * 1.6)
                source_card(s, T.MARGIN, yy, FULL, Inches(1.48), d)
                shapes.append(s.shapes[-1])
            T.animate(s, shapes, step=150)

    # ================================================== 03 GRAPH
    section(prs, 3, "Building the graph",
            "Turning 96 unrelated CSV files into one connected structure — "
            "and knowing how much to trust each connection.")

    s, y = new(prs, "step 1 of 3", "Normalisation")
    T.text(s, T.MARGIN, y, FULL, Inches(0.5),
           "Before anything can be matched, text has to be reduced to a "
           "comparable form. Four functions do this, and every loader uses "
           "them.", size=13.5, color=T.MUTED)
    T.table(s, T.MARGIN, y + Inches(0.6), FULL,
            ["function", "what it removes", "example"],
            [["fold()", "case, punctuation, extra spaces; hyphens → spaces",
              "ATORVASTATIN CALCIUM → atorvastatin calcium"],
             ["strip_salts()", "salt and hydrate suffixes",
              "atorvastatin calcium trihydrate → atorvastatin"],
             ["strip_stereo()", "stereochemical prefixes",
              "R-salbutamol → salbutamol"],
             ["norm_company()", "Inc, Ltd, GmbH, Pharmaceuticals",
              "Pfizer Inc. → pfizer"]],
            col_w=[2.2, 3.6, 5.0], size=11.5, row_h=Inches(0.46))
    w = T.rect(s, T.MARGIN, y + Inches(3.05), FULL, Inches(0.95), T.PANEL_2,
               radius=0.05)
    T.solid(T.rect(s, T.MARGIN, y + Inches(3.05), Pt(3.5), Inches(0.95)),
            T.AMBER)
    T.text(s, T.MARGIN + Inches(0.32), y + Inches(3.2), FULL - Inches(0.7),
           Inches(0.7),
           "These are not cosmetic. fold() turning hyphens into spaces once "
           "broke a filter that tested for \"0-unassigned\" — by the time the "
           "value arrived it was \"0 unassigned\", and the filter silently "
           "matched nothing.", size=12, color=T.MUTED, line=1.3)
    T.animate(s, [w], step=150)

    s, y = new(prs, "step 2 of 3", "The resolver: four tiers, each recording itself")
    T.table(s, T.MARGIN, y, FULL,
            ["", "tier", "matches on", "written to the edge", "confidence"],
            [["1", "exact", "fold(name) → UNII", "unii", "highest"],
             ["2", "salt", "strip_salts(name) → UNII", "salt", "high"],
             ["3", "stereo", "strip_stereo(name) → UNII", "stereo",
              "high, guarded"],
             ["4", "alias", "brand names, research codes", "synonym", "good"],
             ["—", "miss", "nothing matched", "provisional", "weak"]],
            col_w=[0.6, 1.6, 4.0, 2.4, 2.2], size=11.5, row_h=Inches(0.42))
    c1 = T.card(s, T.MARGIN, y + Inches(2.85), Inches(5.85), Inches(1.45),
                "Every edge carries its evidence",
                "match_method travels with the relationship, so an answer can "
                "be filtered by how it was established. Nothing else in this "
                "space exposes that.", T.GREEN, title_size=14, body_size=11.5)
    c2 = T.card(s, T.MARGIN + Inches(6.05), y + Inches(2.85), Inches(5.85),
                Inches(1.45), "First writer wins",
                "Loading order is an authority ranking: GSRS preferred names, "
                "then synonyms, then ChEMBL. Collisions are recorded, never "
                "overwritten.", T.BLUE, title_size=14, body_size=11.5)
    T.animate(s, [c1, c2], step=170)

    s, y = new(prs, "step 3 of 3", "Two failures that shaped the design")
    a = T.card(s, T.MARGIN, y, Inches(5.85), Inches(3.4),
               "The deferred stereo tier",
               "The table must hold the plain form, because the prefix is "
               "usually on the query, not the registered name — GSRS lists "
               "\"Salbutamol\", a source writes \"R-Salbutamol\".\n\n"
               "But that same entry must not let \"Levo-cetirizine\" reach "
               "cetirizine's identifier: they are different drugs. Whether it "
               "would depends on data not yet loaded.\n\n"
               "So the tier is built only after every name is known — propose "
               "an entry for each, then withdraw any that could bridge two "
               "substances.", T.BLUE, title_size=15, body_size=11.5)
    b = T.card(s, T.MARGIN + Inches(6.05), y, Inches(5.85), Inches(3.4),
               "The merge that went wrong",
               "Unresolved names still get a node, so a real product is never "
               "discarded for having an unusual ingredient string.\n\n"
               "Those provisional nodes were originally merged into ChEMBL by "
               "name. ChEMBL contains descriptive names — so \"platinum "
               "complex\" matched 248 distinct molecules and absorbed them "
               "into one node. Across the build, 22,125 provisional nodes "
               "swallowed 45,891 molecules.\n\n"
               "Caught by the validator's fan-out check, not in production.",
               T.AMBER, title_size=15, body_size=11.5)
    T.animate(s, [a, b], step=200)

    s, y = new(prs, "identity", "Every node key is globally unique")
    T.text(s, T.MARGIN, y, FULL, Inches(0.45),
           "The bulk importer resolves every endpoint in a single ID space — "
           "so two labels sharing a key silently become one node.",
           size=13.5, color=T.MUTED)
    T.table(s, T.MARGIN, y + Inches(0.55), Inches(7.3),
            ["namespace", "entity", "value"],
            [["UNII:", "Substance", "FDA unique ingredient identifier"],
             ["CHEMBL:", "Substance", "used when no UNII exists"],
             ["NCT: EUCTR: CTIS:", "ClinicalTrial", "registry's own id"],
             ["FDA: EMA: MHRA: CA:", "Product", "agency's own product id"],
             ["MESH:", "Disease", "MeSH descriptor"],
             ["ID:", "Identifier", "prefixed — see right"]],
            col_w=[2.4, 2.0, 3.4], size=11, row_h=Inches(0.38))
    warn = T.rect(s, T.MARGIN + Inches(7.55), y + Inches(0.55), Inches(4.35),
                  Inches(3.0), T.PANEL_2, radius=0.05)
    T.solid(T.rect(s, T.MARGIN + Inches(7.55), y + Inches(0.55), Pt(3.5),
                   Inches(3.0)), T.RED)
    T.text(s, T.MARGIN + Inches(7.85), y + Inches(0.74), Inches(3.75),
           Inches(2.7),
           [("A real collision, caught before production",
             {"size": 13, "color": T.TEXT, "bold": True, "space_after": 6}),
            ("MESH:D000544 was the key of a Disease and also of an Identifier "
             "recording that same MeSH code. On import they would have merged "
             "into one node that was half disease, half identifier.\n\n"
             "Every identifier is now created through a single function that "
             "prefixes its key. Nothing else may create one.",
             {"size": 11.5, "color": T.MUTED})], line=1.3)
    T.animate(s, [warn], step=150)

    s, y = new(prs, "execution", "Load order is a dependency chain")
    T.text(s, T.MARGIN, y, FULL, Inches(0.4),
           "Each stage fills the lookups the next one reads. Running them out "
           "of order does not fail — it silently produces fewer edges.",
           size=13, color=T.MUTED)
    steps = [("1–3", "Vocabularies + substance spine",
              "ATC tree, then GSRS builds the resolver, then the stereo tier "
              "is finalised"),
             ("4–7", "ChEMBL + reference",
              "2.9M molecules, synonyms, InChIKey structures, salt hierarchy"),
             ("8–9", "Targets + mechanisms",
              "UniProt-keyed proteins; one file yields both TARGETS and "
              "HAS_MECHANISM"),
             ("10", "Disease",
              "MeSH spine and the MeSH/EFO crosswalk — must precede any "
              "disease text matching"),
             ("11–12", "Products + trials",
              "Ten agencies, then ten registries with WHO last"),
             ("13–15", "Safety, variants, literature",
              "Aggregated FAERS, filtered ClinVar, title-matched publications")]
    shapes = []
    for i, (num, head, body) in enumerate(steps):
        yy = y + Inches(0.55 + i * 0.66)
        r = T.rect(s, T.MARGIN, yy, FULL, Inches(0.58), T.PANEL, radius=0.09)
        T.text(s, T.MARGIN + Inches(0.25), yy + Inches(0.14), Inches(0.9),
               Inches(0.3), num, size=13, color=T.BLUE, bold=True)
        T.text(s, T.MARGIN + Inches(1.15), yy + Inches(0.14), Inches(3.4),
               Inches(0.3), head, size=12.5, color=T.TEXT, bold=True)
        T.text(s, T.MARGIN + Inches(4.7), yy + Inches(0.16), Inches(7.0),
               Inches(0.3), body, size=11, color=T.MUTED)
        shapes.append(r)
    T.animate(s, shapes, step=110)

    s, y = new(prs, "the result", "22 entity types, 32 relationship types")
    img = HERE.parent / "graph" / "schema_phase2.png"
    if img.exists():
        pic = s.shapes.add_picture(str(img), T.MARGIN + Inches(0.4), y,
                                   width=Inches(11.1))
        if pic.height > Inches(4.5):
            ratio = Inches(4.5) / pic.height
            pic.height, pic.width = Inches(4.5), int(pic.width * ratio)
            pic.left = int((T.W - pic.width) / 2)
        T.animate(s, [pic], step=150)

    _entities(prs, new)
    _edges(prs, new)

    s, y = new(prs, "trust", "Not all edges are equal — and the graph says so")
    T.table(s, T.MARGIN, y, FULL,
            ["match_method", "meaning", "how to treat it"],
            [["structured", "the source stated the relationship outright",
              "high confidence"],
             ["unii / salt / stereo", "resolved through an identifier tier",
              "high confidence"],
             ["symbol", "joined on a gene symbol", "good"],
             ["derived", "computed from the file's own agency or location",
              "high, but our inference"],
             ["aggregated", "counted across many reports",
              "sound in aggregate, not per report"],
             ["name", "free prose matched against a dictionary",
              "a hint — good for counting, not for citing"],
             ["provisional", "the name never resolved", "weak"]],
            col_w=[2.6, 5.2, 4.0], size=11.5, row_h=Inches(0.4))
    n = T.rect(s, T.MARGIN, y + Inches(3.4), FULL, Inches(0.85), T.PANEL_2,
               radius=0.05)
    T.text(s, T.MARGIN + Inches(0.32), y + Inches(3.56), FULL - Inches(0.7),
           Inches(0.6),
           "This is the difference between a demo and an auditable system. An "
           "answer can be filtered to structured evidence only, and the "
           "provenance of every row is recorded alongside it.",
           size=12.5, color=T.MUTED, line=1.3)
    T.animate(s, [n], step=150)

    s, y = new(prs, "quality gate", "Nothing reaches production unvalidated")
    T.code(s, T.MARGIN, y, Inches(6.0), Inches(2.5),
           "referential integrity   every endpoint exists\n"
           "key uniqueness          no key used by two labels\n"
           "resolution quality      provisional share in bounds\n"
           "fan-out outliers        no node absorbing too much\n"
           "source coverage         every declared file read\n"
           "biology fixtures        atorvastatin -> HMG-CoA\n"
           "                        pembrolizumab -> PD-1\n"
           "                        erenumab -> CGRP receptor")
    T.text(s, T.MARGIN + Inches(6.4), y, Inches(5.5), Inches(2.6),
           bullets([
               ("Why this matters commercially",
                "A bulk import replaces the store outright, with no "
                "transaction. An unchecked build silently becomes the live "
                "graph — and the failure mode is a confident wrong answer, "
                "not an error."),
               ("It has earned its place",
                "The validator caught six structural defects, including the "
                "key collision and the platinum-complex merge, before either "
                "reached a user.")]), line=1.3)
    T.animate(s, [], step=150)

    # ================================================== 04 EMBEDDING
    section(prs, 4, "Embedding entities",
            "Why a knowledge graph still needs vectors — and only for certain "
            "concepts.", T.VIOLET)

    s, y = new(prs, "the gap", "Exact search cannot bridge these", T.VIOLET)
    T.text(s, T.MARGIN, y, FULL, Inches(0.4),
           "Full-text handles names well. It cannot connect a query to a "
           "concept it shares no characters with.", size=13.5, color=T.MUTED)
    T.table(s, T.MARGIN, y + Inches(0.55), FULL,
            ["question arrives as", "full-text finds", "the right node",
             "vector score"],
            [["NSCLC", "nothing", "Non-small cell lung cancer", "0.832"],
             ["statins", "nothing", "HMG-CoA reductase inhibitor", "0.737"],
             ["COPD", "nothing", "Chronic obstructive pulmonary disease",
              "0.803"],
             ["heart problems", "nothing", "heart disorder", "0.873"]],
            col_w=[2.6, 2.0, 5.0, 1.8], size=12, head_color=T.VIOLET,
            row_h=Inches(0.44))
    n = T.rect(s, T.MARGIN, y + Inches(2.9), FULL, Inches(1.2), T.PANEL,
               radius=0.05)
    T.solid(T.rect(s, T.MARGIN, y + Inches(2.9), Pt(3.5), Inches(1.2)),
            T.VIOLET)
    T.text(s, T.MARGIN + Inches(0.32), y + Inches(3.08), FULL - Inches(0.7),
           Inches(0.9),
           [("Every one of these returned nothing before embedding.",
             {"size": 14, "color": T.TEXT, "bold": True, "space_after": 4}),
            ("A researcher does not type \"Carcinoma, Non-Small-Cell Lung\". "
             "They type NSCLC. Without this layer the graph is only reachable "
             "by people who already know its vocabulary.",
             {"size": 12, "color": T.MUTED})])
    T.animate(s, [n], step=150)

    s, y = new(prs, "scope", "Which entities are embedded, and which are not",
               T.VIOLET)
    T.table(s, T.MARGIN, y, Inches(6.6),
            ["embedded", "nodes", "why"],
            [["Disease", "24,488", "arrives as acronyms and lay phrasing"],
             ["DrugClass", "6,996", "\"statins\" is nobody's ATC name"],
             ["AdverseEvent", "6,981", "MedDRA is clinical; complaints are not"],
             ["Mechanism", "1,967", "described in prose that varies freely"]],
            col_w=[2.2, 1.5, 4.4], size=11.5, head_color=T.VIOLET,
            row_h=Inches(0.42))
    T.text(s, T.MARGIN, y + Inches(2.3), Inches(6.6), Inches(0.4),
           "40,432 nodes embedded in total.", size=12.5, color=T.VIOLET,
           bold=True)
    c = T.card(s, T.MARGIN + Inches(6.9), y, Inches(5.0), Inches(2.9),
               "Not embedded, on purpose",
               "Substance, Product, Target, Company and ClinicalTrial are "
               "matched by identifier or exact name — where full-text already "
               "wins, and a vector would introduce plausible wrong answers.\n\n"
               "Embedding 3.07M substances would also cost far more than it "
               "returns: 93% of them carry no name at all, only an "
               "identifier.", T.BLUE, title_size=14, body_size=11.5)
    T.animate(s, [c], step=150)

    s, y = new(prs, "model selection", "Measured, not assumed", T.VIOLET)
    T.text(s, T.MARGIN, y, FULL, Inches(0.4),
           "SapBERT was compared against bge-m3 — the model already in use for "
           "documents — on the same probes.", size=13.5, color=T.MUTED)
    T.table(s, T.MARGIN, y + Inches(0.55), Inches(6.6),
            ["", "SapBERT", "bge-m3"],
            [["probes correct", "6 / 6", "5 / 6"],
             ["score on correct answers", "0.74 – 0.87", "0.43 – 0.61"],
             ["dimensions", "768", "1024"]],
            col_w=[3.0, 2.0, 2.0], size=12, head_color=T.VIOLET,
            row_h=Inches(0.44))
    c = T.card(s, T.MARGIN + Inches(6.9), y + Inches(0.55), Inches(5.0),
               Inches(2.2), "The range matters more than the hit rate",
               "SapBERT is trained on biomedical entity linking — precisely "
               "this task. Its correct answers sit far from its wrong ones, "
               "so a rejection threshold is meaningful. We reject below 0.60.\n\n"
               "bge-m3's correct answers score close to its incorrect ones, "
               "leaving no way to tell them apart.", T.VIOLET, title_size=14,
               body_size=11.5)
    n2 = T.rect(s, T.MARGIN, y + Inches(2.95), FULL, Inches(1.0), T.PANEL_2,
                radius=0.05)
    T.text(s, T.MARGIN + Inches(0.32), y + Inches(3.12), FULL - Inches(0.7),
           Inches(0.75),
           "Both models are kept, on purpose. Documents stay on bge-m3: "
           "passage retrieval is a different task from entity linking, and "
           "re-embedding 3.24M chunks to unify them would cost GPU hours to "
           "make retrieval worse.", size=12, color=T.MUTED, line=1.3)
    T.animate(s, [c, n2], step=170)

    # ================================================== 05 VECTOR STORE
    section(prs, 5, "The document store",
            "93,505 regulatory documents, chunked, embedded, and joined back "
            "to the graph.", T.VIOLET)

    s, y = new(prs, "corpus", "What is in the document store", T.VIOLET)
    T.table(s, T.MARGIN, y, FULL,
            ["source", "chunks", "what these documents are"],
            [["MHRA (UK)", "2,168,590",
              "Public Assessment Reports, SPCs, patient leaflets"],
             ["EMA (Europe)", "988,041",
              "EPARs — assessment history, SmPC, package leaflet"],
             ["PMDA (Japan)", "79,471", "Review reports and package inserts"],
             ["DHA (Dubai)", "2,940", "Circulars and product notices"],
             ["Oman MoH", "806", "Regulatory notices"],
             ["Abu Dhabi DoH", "782", "Circulars and product notices"],
             ["LOINC / NHRA / ANZCTR / Qatar", "126", "Standards and notices"]],
            col_w=[3.0, 2.0, 7.0], size=11.5, head_color=T.VIOLET,
            row_h=Inches(0.38))
    n = T.rect(s, T.MARGIN, y + Inches(3.25), FULL, Inches(0.95), T.PANEL,
               radius=0.05)
    T.solid(T.rect(s, T.MARGIN, y + Inches(3.25), Pt(3.5), Inches(0.95)),
            T.VIOLET)
    T.text(s, T.MARGIN + Inches(0.32), y + Inches(3.42), FULL - Inches(0.7),
           Inches(0.7),
           "3,240,756 chunks across 92,397 documents. The Gulf sources are "
           "small in volume and disproportionate in value — they are the only "
           "machine-readable coverage of those markets we have found.",
           size=12.5, color=T.MUTED, line=1.3)
    T.animate(s, [n], step=150)

    s, y = new(prs, "chunking", "How documents are split, and why", T.VIOLET)
    T.text(s, T.MARGIN, y, Inches(6.0), Inches(2.9),
           bullets([
               ("Heading-aware, not fixed-size",
                "Chunks follow the document's own section structure. A "
                "regulatory document is already organised — contraindications, "
                "dosing, warnings — and cutting every 500 tokens splits a "
                "contraindication in half."),
               ("Why that matters for retrieval",
                "A question about warnings should return the warnings "
                "section, not its last paragraph plus the start of the next "
                "one."),
               ("Deterministic chunk IDs",
                "Re-ingesting a known document overwrites in place instead of "
                "duplicating. Re-running is safe.")]), line=1.3)
    T.code(s, T.MARGIN + Inches(6.4), y, Inches(5.5), Inches(1.9),
           "chunk payload\n"
           "  source     ema.europa.eu\n"
           "  s3_key     ..._PL347710263_par.pdf\n"
           "  doc_type   spc | pil | par\n"
           "  heading    4.3 Contraindications\n"
           "  molecule_id  (reserved)")
    c = T.card(s, T.MARGIN + Inches(6.4), y + Inches(2.1), Inches(5.5),
               Inches(1.15), "Incremental by design",
               "Ingest compares each object's S3 ETag against what is "
               "indexed. A recent run skipped 92,397 unchanged documents and "
               "examined 1,038.", T.GREEN, title_size=13.5, body_size=11)
    T.animate(s, [c], step=150)

    s, y = new(prs, "infrastructure", "Building the first snapshot", T.VIOLET)
    T.text(s, T.MARGIN, y, FULL, Inches(0.4),
           "Embedding 93,505 documents is a one-time cost with a very "
           "different shape from keeping them current.", size=13.5,
           color=T.MUTED)
    a = T.card(s, T.MARGIN, y + Inches(0.55), Inches(5.85), Inches(1.8),
               "First snapshot — GPU",
               "The initial corpus was embedded on a rented RunPod GPU "
               "instance. On CPU the same work is roughly seven minutes per "
               "document — days of wall clock.", T.VIOLET, title_size=14,
               body_size=11.5)
    b = T.card(s, T.MARGIN + Inches(6.05), y + Inches(0.55), Inches(5.85),
               Inches(1.8), "Steady state — CPU",
               "Weekly deltas are tens of documents, which the always-on host "
               "absorbs without a GPU. We rent capacity for backfills and own "
               "it for updates.", T.GREEN, title_size=14, body_size=11.5)
    for i, (v, l, col) in enumerate([("3,240,756", "chunks live", T.VIOLET),
                                     ("92,397", "documents indexed", T.VIOLET),
                                     ("1,024", "dimensions (bge-m3)", T.BLUE),
                                     ("~27 GB", "vectors on disk", T.MUTED)]):
        T.stat(s, T.MARGIN + Inches(i * 3.0), y + Inches(2.6), Inches(2.8),
               v, l, col, size=28, lsize=10)
    T.animate(s, [a, b], step=180)

    # ================================================== 06 ORCHESTRATION
    section(prs, 6, "Orchestration",
            "The part that makes this a product rather than a dataset: it "
            "stays current without anyone touching it.", T.GREEN)

    s, y = new(prs, "scheduling", "One clock, and everything else reacts",
               T.GREEN)
    _dag(s, y)

    s, y = new(prs, "synchronisation", "Why a sync can never race a scrape",
               T.GREEN)
    T.text(s, T.MARGIN, y, Inches(6.1), Inches(3.0),
           bullets([
               ("Event-driven, not scheduled",
                "The sync pipelines have no clock of their own. They wake when "
                "data actually lands in S3, so they cannot run against a "
                "half-written snapshot."),
               ("Routed by content",
                "A source that publishes only CSVs wakes the graph. A source "
                "that publishes documents wakes the vector store. Three "
                "sources publish both and wake both."),
               ("Derived, not maintained",
                "The graph's trigger list is computed from the same file that "
                "declares what it reads — so adding a source updates the "
                "trigger with nothing to keep in step by hand."),
               ("A quiet week costs nothing",
                "No publish, no build, no spend.")]), line=1.28)
    proof = T.rect(s, T.MARGIN + Inches(6.5), y, Inches(5.4), Inches(3.1),
                   T.PANEL, radius=0.05)
    T.text(s, T.MARGIN + Inches(6.8), y + Inches(0.22), Inches(4.8),
           Inches(0.3), "Verified end to end", size=14, color=T.GREEN,
           bold=True)
    T.text(s, T.MARGIN + Inches(6.8), y + Inches(0.65), Inches(4.8),
           Inches(2.2),
           [("One scraper was enabled and nothing else was touched:",
             {"size": 11.5, "color": T.MUTED, "space_after": 8}),
            ("scrape  +4.6 MB of new trial data",
             {"size": 11.5, "color": T.TEXT, "font": T.FONT_MONO,
              "space_after": 4}),
            ("trigger  dataset commit woke the graph",
             {"size": 11.5, "color": T.TEXT, "font": T.FONT_MONO,
              "space_after": 4}),
            ("build    24 min, driven over SSH",
             {"size": 11.5, "color": T.TEXT, "font": T.FONT_MONO,
              "space_after": 4}),
            ("validate 0 failures",
             {"size": 11.5, "color": T.TEXT, "font": T.FONT_MONO,
              "space_after": 4}),
            ("import   8 min, store replaced",
             {"size": 11.5, "color": T.TEXT, "font": T.FONT_MONO,
              "space_after": 8}),
            ("Result: +860 trials, appearing identically in both the "
             "registry-specific and total counts.",
             {"size": 11.5, "color": T.GREEN})], line=1.3)
    T.animate(s, [proof], step=150)

    # ================================================== 07 LIMITS
    section(prs, 7, "Limits and roadmap",
            "What the current hardware costs us — stated plainly, because it "
            "is the main thing standing between today and the next step.",
            T.AMBER)

    s, y = new(prs, "constraint", "Running a distributed system on two small "
               "machines", T.AMBER)
    cards = []
    for i, (t, b) in enumerate([
            ("Two hosts is a floor, not a choice",
             "Qdrant and Neo4j both want memory resident. On one 16 GB box "
             "the larger evicts the smaller, and a swapped page cache is "
             "slower than a small one."),
            ("Cross-cloud by accident",
             "The graph runs on Azure and everything else on AWS. Each build "
             "moves ~15 GB across that boundary, and Azure cannot use an IAM "
             "instance role — so that host needs key files."),
            ("Orchestration reaches across hosts",
             "Airflow drives the graph host over SSH. It works and is "
             "proven, but it is more moving parts than one machine would "
             "need."),
            ("Two vCPU is the query ceiling",
             "Point lookups are 21 ms and two-hop traversals 13–38 ms, but "
             "roughly two concurrent queries run at full speed. This is a "
             "demo-scale box.")]):
        col, row = i % 2, i // 2
        cards.append(T.card(s, T.MARGIN + col * Inches(6.05),
                            y + row * Inches(1.7), Inches(5.85),
                            Inches(1.55), t, b, T.AMBER, title_size=14,
                            body_size=11.5))
    T.animate(s, cards, step=160)

    s, y = new(prs, "what it costs today", "Three things the hardware blocks",
               T.AMBER)
    T.table(s, T.MARGIN, y, FULL,
            ["blocked", "what it would add", "what it needs"],
            [["OpenAlex literature",
              "8.7M works already downloaded — 1,375× current publication "
              "coverage", "~10 GB disk on the graph host"],
             ["Concurrent users",
              "More than ~2 simultaneous analyst queries at full speed",
              "More vCPU"],
             ["1,038 scanned documents",
              "Image-only PDFs that currently produce zero chunks",
              "OCR (tesseract), not yet installed"]],
            col_w=[3.0, 5.5, 3.5], size=11.5, row_h=Inches(0.52),
            head_color=T.AMBER)
    n = T.rect(s, T.MARGIN, y + Inches(2.6), FULL, Inches(1.5), T.PANEL,
               radius=0.05)
    T.solid(T.rect(s, T.MARGIN, y + Inches(2.6), Pt(3.5), Inches(1.5)),
            T.GREEN)
    T.text(s, T.MARGIN + Inches(0.32), y + Inches(2.8), FULL - Inches(0.7),
           Inches(1.15),
           [("None of these are design problems.",
             {"size": 15, "color": T.TEXT, "bold": True, "space_after": 5}),
            ("Every one is a line item. The architecture already separates "
             "storage from compute and code from host — the graph host holds "
             "no irreplaceable state and can be rebuilt from S3 in about "
             "twenty-five minutes. Scaling is provisioning, not rewriting.",
             {"size": 12.5, "color": T.MUTED})], line=1.32)
    T.animate(s, [n], step=150)

    s, y = new(prs, "what exists today", "Where the platform stands", T.GREEN)
    for i, (v, l, col) in enumerate([
            (C.millions(C.live_counts()["_nodes"]), "graph nodes", T.BLUE),
            (C.millions(C.live_counts()["_edges"]), "relationships", T.BLUE),
            ("3.24M", "document chunks", T.VIOLET),
            ("41", "live sources", T.CYAN)]):
        T.stat(s, T.MARGIN + Inches(i * 3.0), y, Inches(2.8), v, l, col,
               size=38, lsize=10.5)
    rows = [["Graph correctness", "0 validation failures on every production build"],
            ["Sync", "Proven end to end — scrape to live graph, unattended"],
            ["Entity retrieval", "Acronyms and lay phrasing resolve at 0.74–0.87"],
            ["Document retrieval", "3.24M chunks, incremental, deterministic"],
            ["Recovery", "Graph rebuilds from S3 in ~25 minutes"],
            ["Documentation", "Architecture, data and query guides generated from the code"]]
    T.table(s, T.MARGIN, y + Inches(1.3), FULL, ["capability", "status"],
            rows, col_w=[3.2, 8.7], size=12, head_color=T.GREEN,
            row_h=Inches(0.42))
    T.animate(s, [], step=150)

    # ---------------------------------------------------------------- end
    s = T.blank(prs)
    T.bg(s)
    _n[0] += 1
    T.solid(T.rect(s, 0, 0, T.W, Inches(0.09)), T.BLUE)
    a = T.text(s, T.MARGIN, Inches(2.6), Inches(11), Inches(1.2),
               "The data was never the moat.\nConnecting it is.",
               size=40, color=T.TEXT, bold=True, line=1.15)
    b = T.text(s, T.MARGIN, Inches(4.3), Inches(10), Inches(1.0),
               "Forty-one sources reconciled into one auditable structure, "
               "kept current automatically, with the evidence for every "
               "connection recorded alongside it.",
               size=15, color=T.MUTED, line=1.4)
    T.footer(s, _n[0])
    T.transition(s, "fade")
    T.animate(s, [a, b], step=220)

    prs.save(OUT)
    print(f"wrote {OUT}  ({OUT.stat().st_size/1024:.0f} KB, {_n[0]} slides)")


# --------------------------------------------------------------------------
def _cat_title(cat, i):
    base = {
        "Drug_Substance_Reference": "The substance spine",
        "Clinical_Trials_Pipeline_Intelligence": "Trial registries",
        "Regulatory_Approvals": "Agency registers",
        "Safety_Pharmacovigilance": "Post-market safety",
        "Targets_Genomics_Biomarkers": "Biology and genomics",
        "Ontologies_Standards": "Controlled vocabularies",
        "Literature_Evidence": "Published evidence",
        "MENA_GCC_Regulatory_Market": "Gulf and MENA coverage",
    }[cat]
    return base if i == 0 else base + " (continued)"


def _architecture(s, y):
    """The system diagram."""
    lake = T.rect(s, T.MARGIN, y, FULL, Inches(0.95), T.PANEL_2, radius=0.05)
    T.solid(T.rect(s, T.MARGIN, y, Pt(3.5), Inches(0.95)), T.CYAN)
    T.text(s, T.MARGIN + Inches(0.3), y + Inches(0.14), Inches(6),
           Inches(0.3), "S3 DATA LAKE  ·  moine-data", size=13,
           color=T.CYAN, bold=True)
    T.text(s, T.MARGIN + Inches(0.3), y + Inches(0.46), Inches(9),
           Inches(0.3),
           "432 CSV files  ·  93,505 documents  ·  50 sources, full "
           "snapshots overwritten each run", size=11.5, color=T.MUTED)
    T.text(s, T.MARGIN + Inches(9.6), y + Inches(0.28), Inches(2.2),
           Inches(0.4), "50 scrapers →", size=12, color=T.DIM,
           align=PP_ALIGN.RIGHT)

    y2 = y + Inches(1.25)
    g = T.rect(s, T.MARGIN, y2, Inches(5.85), Inches(2.55), T.PANEL,
               radius=0.05)
    T.solid(T.rect(s, T.MARGIN, y2, Pt(3.5), Inches(2.55)), T.BLUE)
    T.text(s, T.MARGIN + Inches(0.3), y2 + Inches(0.16), Inches(5),
           Inches(0.3), "KNOWLEDGE GRAPH  ·  Azure host", size=12.5,
           color=T.BLUE, bold=True)
    T.text(s, T.MARGIN + Inches(0.3), y2 + Inches(0.52), Inches(5.2),
           Inches(1.9),
           [("35 sources  →  96 CSV files",
             {"size": 11.5, "color": T.MUTED, "space_after": 6}),
            ("build → validate → import  (Neo4j)",
             {"size": 11.5, "color": T.MUTED, "space_after": 10}),
            (f"{C.live_counts()['_nodes']:,} nodes   "
             f"{C.live_counts()['_edges']:,} relationships",
             {"size": 13, "color": T.TEXT, "bold": True, "space_after": 8}),
            ("22 entity types · 32 relationship types · 3 full-text indexes",
             {"size": 11, "color": T.DIM, "space_after": 6}),
            ("40,432 entity embeddings (SapBERT)",
             {"size": 11, "color": T.DIM})], line=1.3)

    v = T.rect(s, T.MARGIN + Inches(6.05), y2, Inches(5.85), Inches(2.55),
               T.PANEL, radius=0.05)
    T.solid(T.rect(s, T.MARGIN + Inches(6.05), y2, Pt(3.5), Inches(2.55)),
            T.VIOLET)
    T.text(s, T.MARGIN + Inches(6.35), y2 + Inches(0.16), Inches(5),
           Inches(0.3), "DOCUMENT STORE  ·  AWS host", size=12.5,
           color=T.VIOLET, bold=True)
    T.text(s, T.MARGIN + Inches(6.35), y2 + Inches(0.52), Inches(5.2),
           Inches(1.9),
           [("10 sources  →  93,505 documents",
             {"size": 11.5, "color": T.MUTED, "space_after": 6}),
            ("extract → chunk → embed → upsert  (Qdrant)",
             {"size": 11.5, "color": T.MUTED, "space_after": 10}),
            ("3,240,756 chunks   92,397 documents",
             {"size": 13, "color": T.TEXT, "bold": True, "space_after": 8}),
            ("bge-m3, 1024 dimensions · incremental by S3 ETag",
             {"size": 11, "color": T.DIM, "space_after": 6}),
            ("Airflow also runs here, driving the graph host over SSH",
             {"size": 11, "color": T.DIM})], line=1.3)

    j = T.rect(s, T.MARGIN, y2 + Inches(2.8), FULL, Inches(0.72), T.PANEL_2,
               radius=0.05)
    T.text(s, T.MARGIN + Inches(0.3), y2 + Inches(2.95), FULL - Inches(0.6),
           Inches(0.5),
           "The two stores are joined by identifiers, not shared storage — "
           "MHRA filenames embed the licence number matching 39,002 graph "
           "identifiers, so graph → identifier → filtered document search is "
           "one path.", size=11.5, color=T.MUTED, line=1.25)
    T.animate(s, [lake, g, v, j], step=200)


def _stages(s, y):
    stages = [
        ("HYDRATE", "restore state", "Bring back the previous CSV and any "
         "cursors. Without this an incremental crawler restarts at page one "
         "every run and never finishes.", T.CYAN),
        ("SCRAPE", "fetch", "Pull from the API, bulk file or pages. Rate "
         "limits, retries and key rotation live here.", T.BLUE),
        ("NORMALISE", "shape", "Into CSV with a stable header, so the graph "
         "declares columns rather than discovering them.", T.BLUE),
        ("STAGE", "write aside", "Write under a per-run prefix. Nothing "
         "downstream can see a half-written file.", T.VIOLET),
        ("COMMIT", "publish", "Copy to the stable key and rewrite "
         "_LATEST.json. This is the atomic moment.", T.GREEN),
        ("EMIT", "signal", "Publish a dataset event. This is what wakes the "
         "graph and the document store.", T.GREEN),
    ]
    shapes = []
    w = Inches(1.87)
    for i, (name, sub, body, col) in enumerate(stages):
        x = T.MARGIN + i * Inches(1.98)
        r = T.rect(s, x, y, w, Inches(2.9), T.PANEL, radius=0.06)
        T.solid(T.rect(s, x, y, w, Pt(3.5)), col)
        T.text(s, x + Inches(0.18), y + Inches(0.2), w - Inches(0.3),
               Inches(0.28), name, size=12, color=col, bold=True)
        T.text(s, x + Inches(0.18), y + Inches(0.48), w - Inches(0.3),
               Inches(0.24), sub, size=9.5, color=T.DIM)
        T.text(s, x + Inches(0.18), y + Inches(0.82), w - Inches(0.34),
               Inches(1.9), body, size=9.5, color=T.MUTED, line=1.25)
        if i < len(stages) - 1:
            T.text(s, x + w + Inches(0.01), y + Inches(1.3), Inches(0.12),
                   Inches(0.3), "›", size=15, color=T.LINE)
        shapes.append(r)
    n = T.rect(s, T.MARGIN, y + Inches(3.15), FULL, Inches(0.85), T.PANEL_2,
               radius=0.05)
    T.text(s, T.MARGIN + Inches(0.32), y + Inches(3.32), FULL - Inches(0.7),
           Inches(0.6),
           "Hydrate exists because of a real failure: an incremental crawler "
           "kept its own cursor, but nothing restored it — so every run "
           "started from page one, hit the timeout and committed nothing.",
           size=12, color=T.MUTED, line=1.3)
    T.animate(s, shapes + [n], step=130)


def _dag(s, y):
    top = T.rect(s, T.MARGIN, y, FULL, Inches(0.8), T.PANEL_2, radius=0.05)
    T.solid(T.rect(s, T.MARGIN, y, Pt(3.5), Inches(0.8)), T.GREEN)
    T.text(s, T.MARGIN + Inches(0.3), y + Inches(0.13), Inches(6),
           Inches(0.3), "scrapers_pipeline   ·   @weekly", size=13,
           color=T.GREEN, bold=True)
    T.text(s, T.MARGIN + Inches(0.3), y + Inches(0.45), Inches(10),
           Inches(0.28),
           "The only thing on a clock. Each source's commit publishes a "
           "dataset event.", size=11.5, color=T.MUTED)

    y2 = y + Inches(1.15)
    a = T.rect(s, T.MARGIN, y2, Inches(5.85), Inches(2.1), T.PANEL,
               radius=0.05)
    T.solid(T.rect(s, T.MARGIN, y2, Pt(3.5), Inches(2.1)), T.BLUE)
    T.text(s, T.MARGIN + Inches(0.3), y2 + Inches(0.16), Inches(5),
           Inches(0.3), "graph_sync", size=13, color=T.BLUE, bold=True)
    T.text(s, T.MARGIN + Inches(0.3), y2 + Inches(0.52), Inches(5.2),
           Inches(1.5),
           [("Wakes on any of 35 CSV-publishing sources",
             {"size": 11.5, "color": T.MUTED, "space_after": 7}),
            ("build (~24 min) → validate → import (~8 min)",
             {"size": 11.5, "color": T.TEXT, "font": T.FONT_MONO,
              "space_after": 7}),
            ("Runs on the Azure host, driven over SSH. Rebuilds the whole "
             "graph rather than patching it.",
             {"size": 11, "color": T.DIM})], line=1.3)

    b = T.rect(s, T.MARGIN + Inches(6.05), y2, Inches(5.85), Inches(2.1),
               T.PANEL, radius=0.05)
    T.solid(T.rect(s, T.MARGIN + Inches(6.05), y2, Pt(3.5), Inches(2.1)),
            T.VIOLET)
    T.text(s, T.MARGIN + Inches(6.35), y2 + Inches(0.16), Inches(5),
           Inches(0.3), "vector_store_sync", size=13, color=T.VIOLET,
           bold=True)
    T.text(s, T.MARGIN + Inches(6.35), y2 + Inches(0.52), Inches(5.2),
           Inches(1.5),
           [("Wakes on any of 10 document-publishing sources",
             {"size": 11.5, "color": T.MUTED, "space_after": 7}),
            ("extract → chunk → embed → upsert",
             {"size": 11.5, "color": T.TEXT, "font": T.FONT_MONO,
              "space_after": 7}),
            ("Incremental by S3 ETag — only changed documents are "
             "re-embedded.", {"size": 11, "color": T.DIM})], line=1.3)

    n = T.rect(s, T.MARGIN, y2 + Inches(2.35), FULL, Inches(0.9), T.PANEL_2,
               radius=0.05)
    T.solid(T.rect(s, T.MARGIN, y2 + Inches(2.35), Pt(3.5), Inches(0.9)),
            T.AMBER)
    T.text(s, T.MARGIN + Inches(0.32), y2 + Inches(2.52), FULL - Inches(0.7),
           Inches(0.65),
           "Both subscribe with OR semantics. A plain list of datasets means "
           "AND in Airflow — the graph would have fired only when all 35 "
           "sources published in the same window, which is to say never, and "
           "silently.", size=12, color=T.MUTED, line=1.3)
    T.animate(s, [top, a, b, n], step=170)


def _entities(prs, new):
    """Entity slides: what each is and where it comes from."""
    items = list(C.NODE_ORIGIN.items())
    per = 5
    for i in range(0, len(items), per):
        part = items[i:i + per]
        s, y = new(prs, "entities" + ("" if i == 0 else " (continued)"),
                   "What each entity is, and where it comes from")
        rows = []
        for label, (src, why) in part:
            rows.append([label, C.NODE_COUNTS.get(label.split(" /")[0], "—"),
                         src, why])
        T.table(s, T.MARGIN, y, FULL,
                ["entity", "nodes", "built from", "how"],
                rows, col_w=[1.7, 1.3, 2.8, 6.1], size=11,
                row_h=Inches(0.72))
        T.animate(s, [], step=150)


def _edges(prs, new):
    """Relationship slides, generated from the same table as the tech doc."""
    order = ["CONTAINS", "DEVELOPS", "APPROVED_BY", "HAS_APPROVAL",
             "PROTECTED_BY", "HAS_EXCLUSIVITY", "IN_CLASS", "TARGETS",
             "HAS_MECHANISM", "INDICATED_FOR", "IS_SALT_OF",
             "SPONSORED_BY", "STUDIES", "TESTED_IN", "CONDUCTED_IN",
             "SAME_STUDY_AS", "HAS_ADVERSE_EVENT", "IN_ORGAN_CLASS",
             "SUBJECT_OF", "VARIANT_IN", "IMPLICATED_IN", "ASSOCIATED_WITH",
             "ABOUT", "MENTIONS", "HAS_IDENTIFIER", "IN_REGION",
             "SUBTYPE_OF", "HAS_ROUTE", "HAS_MODALITY", "BIOSIMILAR_OF",
             "APPROVED_IN", "ISSUED_BY"]
    seen = [k for k in order if k in C.EDGES]
    seen += [k for k in C.EDGES if k not in seen]
    per = 8
    for i in range(0, len(seen), per):
        part = seen[i:i + per]
        s, y = new(prs, "relationships" + ("" if i == 0 else " (continued)"),
                   "How every connection is made")
        rows = []
        for k in part:
            frm, to, method, meaning = C.EDGES[k]
            rows.append([k, f"{frm} → {to}", method,
                         meaning[:96]])
        T.table(s, T.MARGIN, y, FULL,
                ["relationship", "connects", "evidence", "meaning"],
                rows, col_w=[2.2, 3.2, 1.7, 4.8], size=10.5,
                row_h=Inches(0.5))
        T.animate(s, [], step=150)


if __name__ == "__main__":
    build()
