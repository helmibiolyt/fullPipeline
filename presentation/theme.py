#!/usr/bin/env python3
"""Design system for the investor deck: palette, layout primitives, motion.

Everything the deck draws goes through here, so the slides stay consistent
without each one restating fonts and colours.

On motion: python-pptx has no animation API, so transitions and entrance
effects are injected as raw DrawingML. That is worth doing carefully - a
malformed timing tree makes PowerPoint show a "repair" prompt on open, which
in front of investors is worse than a deck with no animation at all. The
timing tree built here follows the shape PowerPoint itself writes.
"""
from __future__ import annotations

import copy

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import _nsmap as NS, qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

# --------------------------------------------------------------------------
# Palette. Dark, high contrast, one accent that carries meaning: blue is the
# graph, violet is the vector store, amber is a caveat.
INK        = RGBColor(0x07, 0x0C, 0x18)   # page
PANEL      = RGBColor(0x11, 0x1A, 0x2B)   # raised surface
PANEL_2    = RGBColor(0x18, 0x24, 0x3A)   # raised surface, lighter
LINE       = RGBColor(0x25, 0x33, 0x4D)
TEXT       = RGBColor(0xEC, 0xF1, 0xF9)
MUTED      = RGBColor(0x93, 0xA6, 0xC4)
DIM        = RGBColor(0x63, 0x76, 0x94)

BLUE       = RGBColor(0x4C, 0x8D, 0xFF)   # graph
CYAN       = RGBColor(0x22, 0xD3, 0xEE)
VIOLET     = RGBColor(0xA7, 0x8B, 0xFA)   # vector store
GREEN      = RGBColor(0x34, 0xD3, 0x99)   # verified
AMBER      = RGBColor(0xFB, 0xBF, 0x24)   # caveat
RED        = RGBColor(0xF8, 0x71, 0x71)

WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

FONT       = "Segoe UI"
FONT_LIGHT = "Segoe UI Light"
FONT_MONO  = "Consolas"

# 16:9
W = Inches(13.333)
H = Inches(7.5)
MARGIN = Inches(0.72)


# --------------------------------------------------------------------------
# Primitives
def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def rect(slide, x, y, w, h, color=None, radius=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(kind, x, y, w, h)
    if radius:
        # adjustment 0..0.5 of the shorter side
        s.adjustments[0] = radius
    s.shadow.inherit = False
    if color is not None:
        solid(s, color)
    else:
        s.fill.background()
        s.line.fill.background()
    return s


def text(slide, x, y, w, h, runs, size=18, color=None, bold=False,
         align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP,
         line=1.25, space_after=0):
    """runs: a string, or a list of (text, {overrides}) tuples per paragraph."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0

    items = [runs] if isinstance(runs, str) else runs
    for i, item in enumerate(items):
        body, over = (item, {}) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = over.get("align", align)
        p.line_spacing = over.get("line", line)
        p.space_after = Pt(over.get("space_after", space_after))
        r = p.add_run()
        r.text = body
        f = r.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.color.rgb = over.get("color", color if color is not None else TEXT)
    return box


def bg(slide, color=INK):
    """Full-bleed background. python-pptx cannot set slide fill directly."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    s.shadow.inherit = False
    solid(s, color)
    return s


def accent_bar(slide, y, color=BLUE, w=Inches(1.15), x=MARGIN, h=Pt(4.5)):
    return solid(rect(slide, x, y, w, h), color)


# Roughly how many characters of the title font fit one line at 34pt across
# the content width. Used to decide whether a title wraps - python-pptx cannot
# measure text, and a title that wraps unnoticed lands on top of the content.
_TITLE_ONE_LINE = 50


def slide_title(slide, title, kicker=None, color=BLUE):
    y = Inches(0.62)
    if kicker:
        text(slide, MARGIN, y, Inches(11), Inches(0.3), kicker.upper(),
             size=12.5, color=color, bold=True)
        y += Inches(0.42)

    two = len(title) > _TITLE_ONE_LINE
    size = 29 if two else 34
    h = Inches(1.25) if two else Inches(0.85)
    text(slide, MARGIN, y, Inches(11.9), h, title,
         size=size, color=TEXT, bold=True, font=FONT, line=1.12)
    bar = Inches(1.12) if two else Inches(0.72)
    accent_bar(slide, y + bar, color)
    return y + bar + Inches(0.33)


def footer(slide, n, label="Biolyt"):
    text(slide, MARGIN, H - Inches(0.52), Inches(6), Inches(0.3), label,
         size=10, color=DIM)
    text(slide, W - MARGIN - Inches(1.2), H - Inches(0.52), Inches(1.2),
         Inches(0.3), str(n), size=10, color=DIM, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title, body, accent=BLUE, title_size=16,
         body_size=12.5, fill=PANEL):
    """A titled panel. The deck's main unit of content."""
    box = rect(slide, x, y, w, h, fill, radius=0.045)
    solid(rect(slide, x, y, Pt(3.5), h), accent)
    pad = Inches(0.28)
    text(slide, x + pad, y + Inches(0.2), w - pad * 2, Inches(0.4), title,
         size=title_size, color=TEXT, bold=True)
    if body:
        text(slide, x + pad, y + Inches(0.2) + Inches(0.34),
             w - pad * 2, h - Inches(0.75), body,
             size=body_size, color=MUTED, line=1.3)
    return box


def stat(slide, x, y, w, value, label, color=BLUE, size=44, lsize=12):
    text(slide, x, y, w, Inches(0.7), value, size=size, color=color, bold=True)
    text(slide, x, y + Inches(0.62), w, Inches(0.4), label.upper(),
         size=lsize, color=MUTED, bold=True)


def table(slide, x, y, w, headers, rows, col_w=None, size=11,
          head_color=BLUE, row_h=Inches(0.32)):
    """A light table: no gridlines, a ruled header, alternating row tint."""
    n = len(rows) + 1
    shape = slide.shapes.add_table(n, len(headers), x, y, w,
                                   row_h * n).table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            shape.columns[i].width = Emu(int(w * cw / total))
    shape.first_row = False
    for j, htxt in enumerate(headers):
        c = shape.cell(0, j)
        c.text = ""
        _cell(c, htxt, size=size, color=head_color, bold=True, fill=INK)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            _cell(shape.cell(i, j), str(val), size=size,
                  color=TEXT if j == 0 else MUTED,
                  fill=PANEL if i % 2 else INK)
    return shape


def _no_border(cell):
    """Remove a cell's borders.

    python-pptx has no API for this and the built-in table styles all draw
    light gridlines, which on a dark slide look like a rendering fault.
    """
    tc = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for el in tc.findall(qn(tag)):
            tc.remove(el)
        ln = tc.makeelement(qn(tag), {})
        ln.append(ln.makeelement(qn("a:noFill"), {}))
        tc.append(ln)


def _cell(c, value, size, color, bold=False, fill=INK):
    _no_border(c)
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.margin_left = c.margin_right = Inches(0.09)
    c.margin_top = c.margin_bottom = Inches(0.03)
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = c.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = value
    r.font.size = Pt(size)
    r.font.name = FONT
    r.font.bold = bold
    r.font.color.rgb = color


def code(slide, x, y, w, h, body, size=11.5, fill=RGBColor(0x0C, 0x14, 0x24)):
    rect(slide, x, y, w, h, fill, radius=0.03)
    solid(rect(slide, x, y, Pt(3), h), CYAN)
    lines = [(l, {"font": FONT_MONO, "size": size,
                  "color": MUTED if not l.startswith("#") else DIM})
             for l in body.split("\n")]
    text(slide, x + Inches(0.24), y + Inches(0.16), w - Inches(0.4),
         h - Inches(0.3), lines, line=1.22)


def chip(slide, x, y, label, color=BLUE, w=None):
    w = w or Inches(0.16 * len(label) + 0.3)
    s = rect(slide, x, y, w, Inches(0.28), color, radius=0.5)
    text(slide, x, y + Inches(0.03), w, Inches(0.22), label, size=10.5,
         color=INK, bold=True, align=PP_ALIGN.CENTER)
    return x + w + Inches(0.1)


def arrow(slide, x, y, w, color=LINE, h=Pt(2)):
    solid(rect(slide, x, y, w, h), color)
    t = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                               x + w, y - Inches(0.045), Inches(0.11),
                               Inches(0.13))
    t.rotation = 90
    t.shadow.inherit = False
    solid(t, color)


# --------------------------------------------------------------------------
# Motion
_TRANS = {
    "fade":  '<p:fade/>',
    "wipe":  '<p:wipe dir="d"/>',
    "push":  '<p:push dir="u"/>',
    "cut":   '<p:cut/>',
}


def transition(slide, kind="fade", ms=500):
    """Slide transition. Written into <p:sld> after <p:cSld>/<p:clrMapOvr>."""
    xml = (f'<p:transition xmlns:p="{NS["p"]}" spd="med" '
           f'advTm="0" advClick="1">{_TRANS[kind]}</p:transition>')
    el = etree.fromstring(xml)
    # p:transition must follow p:clrMapOvr in the schema.
    slide.element.append(el)


_ANIM_HEAD = """<p:timing xmlns:p="{p}" xmlns:a="{a}">
 <p:tnLst>
  <p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
   <p:childTnLst>
    <p:seq concurrent="1" nextAc="seek">
     <p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>
      <p:par><p:cTn id="3" fill="hold">
       <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
       <p:childTnLst>
        <p:par><p:cTn id="4" fill="hold">
         <p:stCondLst><p:cond delay="0"/></p:stCondLst>
         <p:childTnLst>{effects}</p:childTnLst>
        </p:cTn></p:par>
       </p:childTnLst>
      </p:cTn></p:par>
     </p:childTnLst></p:cTn>
     <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
     <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
    </p:seq>
   </p:childTnLst>
  </p:cTn></p:par>
 </p:tnLst>
</p:timing>"""

_EFFECT = """<p:par><p:cTn id="{i0}" presetID="10" presetClass="entr"
   presetSubtype="0" fill="hold" grpId="0" nodeType="{node}">
  <p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>
  <p:childTnLst>
   <p:set>
    <p:cBhvr>
     <p:cTn id="{i1}" dur="1" fill="hold">
      <p:stCondLst><p:cond delay="0"/></p:stCondLst>
     </p:cTn>
     <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
     <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
    </p:cBhvr>
    <p:to><p:strVal val="visible"/></p:to>
   </p:set>
   <p:animEffect transition="in" filter="fade">
    <p:cBhvr>
     <p:cTn id="{i2}" dur="{dur}"/>
     <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
    </p:cBhvr>
   </p:animEffect>
   <p:anim calcmode="lin" valueType="num">
    <p:cBhvr additive="base">
     <p:cTn id="{i3}" dur="{dur}" fill="hold"/>
     <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
     <p:attrNameLst><p:attrName>ppt_y</p:attrName></p:attrNameLst>
    </p:cBhvr>
    <p:tavLst>
     <p:tav tm="0"><p:val><p:strVal val="#ppt_y+.035"/></p:val></p:tav>
     <p:tav tm="100000"><p:val><p:strVal val="#ppt_y"/></p:val></p:tav>
    </p:tavLst>
   </p:anim>
  </p:childTnLst>
 </p:cTn></p:par>"""


def animate(slide, shapes, step=180, dur=420, start=120):
    """Fade-and-rise the given shapes in order, automatically on slide entry.

    The first effect is nodeType="afterEffect" rather than "clickEffect", so
    the sequence plays on its own - a presenter should not have to click
    twelve times to assemble one slide.
    """
    if not shapes:
        return
    effects, uid = [], 5
    for k, sh in enumerate(shapes):
        effects.append(_EFFECT.format(
            i0=uid, i1=uid + 1, i2=uid + 2, i3=uid + 3,
            spid=sh.shape_id, delay=start + k * step, dur=dur,
            node="afterEffect"))
        uid += 4
    xml = _ANIM_HEAD.format(p=NS["p"], a=NS["a"],
                            effects="".join(effects))
    el = etree.fromstring(xml)
    old = slide.element.find(qn("p:timing"))
    if old is not None:
        slide.element.remove(old)
    slide.element.append(el)


def blank(prs):
    """A slide with no placeholders at all - layout 6 in the default master."""
    return prs.slides.add_slide(prs.slide_layouts[6])
