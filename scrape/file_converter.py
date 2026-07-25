"""
Shared file-to-CSV converter using MiniMax M3 model.

Extracts raw text from files (PDF, Word, PowerPoint, Excel),
then uses MiniMax LLM to intelligently structure the content into CSV.
"""

import os
import io
import csv
import json
import logging
import base64
import time
from pathlib import Path
from typing import Optional

import requests

log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# MiniMax API config
# ---------------------------------------------------------------------------
MINIMAX_API_KEY = os.environ.get("MINIMAX_API_KEY", "")
MINIMAX_API_URL = "https://api.minimax.io/v1/text/chatcompletion_v2"
MINIMAX_MODEL = "MiniMax-M3"
MAX_TEXT_CHARS = 200_000  # truncate very large documents


def _load_api_key() -> str:
    """Load API key from env or .env file."""
    key = MINIMAX_API_KEY
    if not key:
        env_path = Path(__file__).parent / ".env"
        if env_path.exists():
            for line in env_path.read_text().splitlines():
                if line.startswith("MINIMAX_API_KEY="):
                    key = line.split("=", 1)[1].strip()
                    break
    return key


# ---------------------------------------------------------------------------
# Text extraction from various file formats
# ---------------------------------------------------------------------------

def extract_text_from_pdf(data: bytes) -> str:
    """Extract all text from a PDF using pdfplumber."""
    import pdfplumber
    text_parts = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(f"--- Page {i+1} ---\n{page_text}")
    return "\n\n".join(text_parts)


def extract_text_from_docx(data: bytes) -> str:
    """Extract text from a Word document."""
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Also extract tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append("\n".join(rows))
    return "\n".join(parts)


def extract_text_from_pptx(data: bytes) -> str:
    """Extract text from a PowerPoint file."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_text.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_text.append(" | ".join(cells))
        if slide_text:
            parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
    return "\n\n".join(parts)


def extract_text_from_excel(data: bytes, ext: str) -> str:
    """Extract text from Excel files."""
    import pandas as pd
    engine = "openpyxl" if ext in (".xlsx", ".xlsm") else "xlrd"
    try:
        xf = pd.ExcelFile(io.BytesIO(data), engine=engine)
    except Exception:
        engine = "xlrd" if engine == "openpyxl" else "openpyxl"
        xf = pd.ExcelFile(io.BytesIO(data), engine=engine)
    parts = []
    for sheet in xf.sheet_names:
        df = pd.read_excel(xf, sheet_name=sheet, header=None, dtype=str)
        df = df.dropna(how="all").fillna("")
        if df.empty:
            continue
        parts.append(f"--- Sheet: {sheet} ---\n{df.to_csv(index=False, header=False)}")
    return "\n\n".join(parts)


def extract_text(data: bytes, ext: str) -> Optional[str]:
    """Extract text from a file based on its extension. Returns None if unsupported."""
    ext = ext.lower()
    try:
        if ext == ".pdf":
            return extract_text_from_pdf(data)
        elif ext in (".docx", ".doc", ".dotx"):
            return extract_text_from_docx(data)
        elif ext in (".pptx", ".ppt"):
            return extract_text_from_pptx(data)
        elif ext in (".xlsx", ".xls", ".xlsm"):
            return extract_text_from_excel(data, ext)
        else:
            return None
    except Exception as e:
        log.error(f"Text extraction failed for {ext}: {e}")
        return None


# ---------------------------------------------------------------------------
# MiniMax LLM call to structure text into CSV
# ---------------------------------------------------------------------------

SYSTEM_PROMPT = """You are a data extraction assistant. You receive raw text extracted from a document (PDF, Word, Excel, or PowerPoint).

Your job:
1. Analyze the text and identify ALL structured/tabular data that can be represented as CSV.
2. If the document contains tables, extract them faithfully as CSV.
3. If the document contains semi-structured data (lists, key-value pairs, records), organize it into a logical CSV format.
4. If the document is purely narrative text (guidelines, policies, circulars) with no extractable structured data, create a CSV with columns: "Section", "Content" — breaking the text into logical sections.
5. If there are multiple distinct tables/datasets, separate them with a line containing only: ===TABLE_SEPARATOR===

Rules:
- Output ONLY raw CSV data, no markdown, no code fences, no explanations.
- Use comma as delimiter, quote fields that contain commas.
- Include a header row for each table.
- Preserve all data — do not summarize or skip rows.
- Clean up formatting artifacts (extra spaces, broken lines) but keep the actual content.
- For Arabic/bilingual content, keep both languages if present."""


def call_minimax(text: str, retries: int = 3) -> Optional[str]:
    """Send text to MiniMax M3 and get structured CSV back."""
    api_key = _load_api_key()
    if not api_key:
        log.error("MINIMAX_API_KEY not set")
        return None

    if len(text) > MAX_TEXT_CHARS:
        log.warning(f"Text too long ({len(text)} chars), truncating to {MAX_TEXT_CHARS}")
        text = text[:MAX_TEXT_CHARS]

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": MINIMAX_MODEL,
        "messages": [
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": f"Extract structured CSV data from this document:\n\n{text}"},
        ],
        "temperature": 0.1,
        "max_tokens": 65536,
    }

    for attempt in range(retries):
        try:
            resp = requests.post(MINIMAX_API_URL, headers=headers, json=payload, timeout=120)
            resp.raise_for_status()
            result = resp.json()
            # OpenAI-compatible response format
            if "choices" in result:
                return result["choices"][0]["message"]["content"]
            log.error(f"Unexpected response format: {list(result.keys())}")
            return None
        except requests.exceptions.RequestException as e:
            log.warning(f"MiniMax API attempt {attempt+1}/{retries} failed: {e}")
            if attempt < retries - 1:
                time.sleep(2 ** attempt)
    return None


def parse_minimax_csv_response(csv_text: str) -> list[list[list[str]]]:
    """
    Parse MiniMax response into list of tables.
    Each table is a list of rows, each row is a list of strings.
    Handles ===TABLE_SEPARATOR=== to split multiple tables.
    """
    # Strip markdown code fences if present
    csv_text = csv_text.strip()
    if csv_text.startswith("```"):
        lines = csv_text.split("\n")
        # Remove first and last ``` lines
        if lines[0].startswith("```"):
            lines = lines[1:]
        if lines and lines[-1].strip() == "```":
            lines = lines[:-1]
        csv_text = "\n".join(lines)

    tables = []
    chunks = csv_text.split("===TABLE_SEPARATOR===")

    for chunk in chunks:
        chunk = chunk.strip()
        if not chunk:
            continue
        try:
            reader = csv.reader(io.StringIO(chunk))
            rows = [row for row in reader if any(cell.strip() for cell in row)]
            if rows:
                tables.append(rows)
        except csv.Error as e:
            log.warning(f"CSV parse error: {e}")
            # Fallback: treat each line as a single-column row
            rows = [[line.strip()] for line in chunk.splitlines() if line.strip()]
            if rows:
                tables.append(rows)

    return tables


# ---------------------------------------------------------------------------
# Main conversion function
# ---------------------------------------------------------------------------

def convert_file_to_csv(data: bytes, ext: str, out_dir: Path, base_name: str) -> list[Path]:
    """
    Convert any supported file to CSV(s) using MiniMax LLM.

    For Excel files that are already structured, uses pandas directly.
    For PDF/Word/PPT, extracts text and sends to MiniMax for structuring.
    For CSV files, just copies them.
    """
    ext = ext.lower()

    # CSV files — just copy
    if ext == ".csv":
        out = out_dir / f"{base_name}.csv"
        out.write_bytes(data)
        return [out]

    # Excel files — already structured, use pandas directly (no LLM needed)
    if ext in (".xlsx", ".xls", ".xlsm"):
        return _excel_to_csvs_direct(data, ext, out_dir, base_name)

    # PDF, Word, PowerPoint — extract text, then LLM
    text = extract_text(data, ext)
    if not text or not text.strip():
        log.warning(f"No text extracted from {base_name}{ext}")
        return []

    log.info(f"Extracted {len(text)} chars from {base_name}{ext}, sending to MiniMax...")
    csv_response = call_minimax(text)
    if not csv_response:
        log.error(f"MiniMax returned no response for {base_name}{ext}")
        # Fallback: save raw text as single-column CSV
        return _text_fallback(text, out_dir, base_name)

    tables = parse_minimax_csv_response(csv_response)
    if not tables:
        log.warning(f"No tables parsed from MiniMax response for {base_name}{ext}")
        return _text_fallback(text, out_dir, base_name)

    paths = []
    for i, rows in enumerate(tables):
        suffix = f"_table{i+1}" if len(tables) > 1 else ""
        out = out_dir / f"{base_name}{suffix}.csv"
        with open(out, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f, quoting=csv.QUOTE_ALL)
            writer.writerows(rows)
        paths.append(out)
        log.info(f"  Wrote {len(rows)} rows to {out.name}")

    return paths


def _excel_to_csvs_direct(data: bytes, ext: str, out_dir: Path, base_name: str) -> list[Path]:
    """Convert Excel to CSV using pandas (no LLM needed for structured data)."""
    import pandas as pd
    engine = "openpyxl" if ext in (".xlsx", ".xlsm") else "xlrd"
    try:
        xf = pd.ExcelFile(io.BytesIO(data), engine=engine)
    except Exception:
        engine = "xlrd" if engine == "openpyxl" else "openpyxl"
        xf = pd.ExcelFile(io.BytesIO(data), engine=engine)

    all_tables = []
    for sheet in xf.sheet_names:
        df = pd.read_excel(xf, sheet_name=sheet, header=None, dtype=str)
        df = df.dropna(how="all").fillna("")
        if df.empty:
            continue
        all_tables.append(df.values.tolist())

    if not all_tables:
        return []

    # Group tables by column count
    from collections import OrderedDict
    groups = OrderedDict()
    for rows in all_tables:
        if not rows:
            continue
        key = len(rows[0])
        if key not in groups:
            groups[key] = []
        groups[key].extend(rows)

    paths = []
    for i, rows in enumerate(groups.values()):
        suffix = f"_topic{i+1}" if len(groups) > 1 else ""
        out = out_dir / f"{base_name}{suffix}.csv"
        with open(out, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f, quoting=csv.QUOTE_ALL)
            writer.writerows(rows)
        paths.append(out)
    return paths


def _text_fallback(text: str, out_dir: Path, base_name: str) -> list[Path]:
    """Fallback: save raw text as single-column CSV."""
    out = out_dir / f"{base_name}_text.csv"
    with open(out, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f, quoting=csv.QUOTE_ALL)
        writer.writerow(["Content"])
        for line in text.splitlines():
            line = line.strip()
            if line:
                writer.writerow([line])
    return [out]
