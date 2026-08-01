"""Text extraction from PDF / DOCX / PPTX, structure-aware where possible.

Returns a list of (page_no, text) blocks. Chunking (chunk.py) then splits these
into section-aware chunks. Kept dependency-light: pymupdf / python-docx /
python-pptx. (Swap in Docling later for richer table/section structure.)
"""
from __future__ import annotations

import os
from pathlib import Path


def extract_blocks(path: str) -> list[tuple[int | None, str]]:
    suf = Path(path).suffix.lower()
    if suf == ".pdf":
        return _pdf(path)
    if suf in (".docx", ".doc"):
        return _docx(path)
    if suf in (".pptx", ".ppt"):
        return _pptx(path)
    return []


#: Below this many characters a page is treated as having no text layer. Not
#: zero: a scanned page often carries a stray header or a page number from the
#: scanner's own stamp, which is enough to look extracted and mean nothing.
_MIN_PAGE_CHARS = 40

#: OCR languages. Arabic matters here - 38 of the 416 SFDA documents are Arabic
#: and a third of those are scans, so an English-only model returns gibberish
#: rather than nothing, which is worse.
OCR_LANGS = os.environ.get("OCR_LANGS", "eng+ara")

#: A document is treated as scanned only when MOST of it has no text layer.
#: The first version decided per page, and that was far too eager: a sparse
#: page is normal in a born-digital report - a table, a section divider, a
#: title page - so a 146-page trial protocol was being OCR'd almost end to
#: end for nothing. It ran for over an hour on one document, invisible,
#: because PyMuPDF calls libtesseract IN-PROCESS: no tesseract process to
#: see, no output until the document finishes.
OCR_PAGE_FRACTION = 0.6

#: And a ceiling regardless, so one pathological file cannot hold the queue.
#: Pages beyond this are left as whatever the text layer gave.
OCR_MAX_PAGES = int(os.environ.get("OCR_MAX_PAGES", "40"))


def _pdf(path: str):
    """Page text, falling back to OCR for pages that have no text layer.

    Without the fallback a scanned PDF extracted to nothing, produced no
    chunks, and was counted by ingest.py as a document successfully ingested.
    36 of the 556 backfilled documents are scans, so 36 SFDA safety alerts
    would have been absent from the corpus with no error anywhere - the same
    silent-absence failure the graph kept producing, in a different system.

    Whether to OCR is decided ONCE per document, not per page. Deciding per
    page looked more careful and was much worse: a sparse page is normal in a
    born-digital report, so a 146-page trial protocol went end-to-end through
    tesseract to recover nothing, and sat on both cores for over an hour with
    no output - PyMuPDF calls libtesseract in-process, so there is no tesseract
    process to notice and the log line only prints once the document is done.
    """
    import fitz  # pymupdf
    out, ocr_pages = [], 0
    with fitz.open(path) as doc:
        pages = [(i, p.get_text("text").strip()) for i, p in enumerate(doc, 1)]

        # Decide ONCE, for the document. Pages with no text are only worth
        # OCRing if the file is a scan; in a document that extracted fine they
        # are blanks, dividers and figure pages, and OCRing them costs minutes
        # to recover nothing.
        blank = sum(1 for _, t in pages if len(t) < _MIN_PAGE_CHARS)
        scanned = pages and blank / len(pages) >= OCR_PAGE_FRACTION

        for i, txt in pages:
            if scanned and len(txt) < _MIN_PAGE_CHARS and ocr_pages < OCR_MAX_PAGES:
                try:
                    page = doc[i - 1]
                    tp = page.get_textpage_ocr(language=OCR_LANGS, full=False,
                                               dpi=200)
                    ocr = page.get_text("text", textpage=tp).strip()
                    if len(ocr) > len(txt):
                        txt = ocr
                        ocr_pages += 1
                except Exception:                              # noqa: BLE001
                    # No tesseract, or a page it cannot handle. Left to the
                    # empty-document check in ingest.py rather than failing.
                    pass
            if txt:
                out.append((i, txt))
    if ocr_pages:
        print(f"    ocr: {ocr_pages} page(s) in {Path(path).name}", flush=True)
    return out


def _docx(path: str):
    from docx import Document
    doc = Document(path)
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    # tables -> pipe-joined rows so tabular content survives
    for t in doc.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                paras.append(" | ".join(cells))
    return [(None, "\n".join(paras))] if paras else []


def _pptx(path: str):
    from pptx import Presentation
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [sh.text for sh in slide.shapes if sh.has_text_frame and sh.text.strip()]
        if texts:
            out.append((i, "\n".join(texts)))
    return out
